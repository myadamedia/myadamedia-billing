import os
import sys

# Ensure site-packages path is included
sys.path.append(r'C:\Users\iwanwAppData\Roaming\Python\Python314\site-packages')
sys.path.append(r'C:\Users\iwanw\AppData\Roaming\Python\Python314\site-packages')

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_pitch_deck(pptx_path):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Color Palette
    DARK_BG = RGBColor(11, 19, 43)       # #0B132B
    CARD_BG = RGBColor(28, 37, 65)       # #1C2541
    ACCENT_CYAN = RGBColor(0, 229, 255)  # #00E5FF
    ACCENT_GOLD = RGBColor(255, 209, 102) # #FFD166
    TEXT_WHITE = RGBColor(255, 255, 255)
    TEXT_MUTED = RGBColor(141, 153, 174) # #8D99AE

    blank_layout = prs.slide_layouts[6]  # Blank slide

    def add_background(slide):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = DARK_BG
        bg.line.fill.background()

    def add_header(slide, number_str, title_text):
        # Header bar
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.733), Inches(0.8))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        
        run_num = p.add_run()
        run_num.text = f"{number_str}  |  "
        run_num.font.bold = True
        run_num.font.size = Pt(14)
        run_num.font.color.rgb = ACCENT_CYAN
        
        run_title = p.add_run()
        run_title.text = title_text.upper()
        run_title.font.bold = True
        run_title.font.size = Pt(22)
        run_title.font.color.rgb = TEXT_WHITE

    # -------------------------------------------------------------
    # SLIDE 1: COVER SLIDE
    # -------------------------------------------------------------
    slide1 = prs.slides.add_slide(blank_layout)
    add_background(slide1)

    # Accent Card Box
    card1 = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(1.5), Inches(10.333), Inches(4.5))
    card1.fill.solid()
    card1.fill.fore_color.rgb = CARD_BG
    card1.line.color.rgb = ACCENT_CYAN
    card1.line.width = Pt(1.5)

    tx_cover = slide1.shapes.add_textbox(Inches(1.8), Inches(1.8), Inches(9.733), Inches(3.9))
    tf1 = tx_cover.text_frame
    tf1.word_wrap = True

    p = tf1.paragraphs[0]
    p.text = "MYADAMEDIA BILLING"
    p.font.bold = True
    p.font.size = Pt(40)
    p.font.color.rgb = ACCENT_CYAN

    p2 = tf1.add_paragraph()
    p2.text = "All-in-One ISP Management & Network Automation Platform"
    p2.font.bold = True
    p2.font.size = Pt(20)
    p2.font.color.rgb = TEXT_WHITE
    p2.space_before = Pt(10)

    p3 = tf1.add_paragraph()
    p3.text = '"Transforming Local ISPs & RTRW Nets into Automated, High-Margin Businesses"'
    p3.font.italic = True
    p3.font.size = Pt(16)
    p3.font.color.rgb = ACCENT_GOLD
    p3.space_before = Pt(15)

    p4 = tf1.add_paragraph()
    p4.text = "Presenter: Management & Technical Team MyAdamedia  |  Contact: admin@myadamedia.com"
    p4.font.size = Pt(13)
    p4.font.color.rgb = TEXT_MUTED
    p4.space_before = Pt(25)

    # -------------------------------------------------------------
    # SLIDE 2: PERMASALAHAN PASAR (THE PROBLEM)
    # -------------------------------------------------------------
    slide2 = prs.slides.add_slide(blank_layout)
    add_background(slide2)
    add_header(slide2, "SLIDE 02", "Permasalahan Utama Pasar (Market Pain Points)")

    problems = [
        ("Billing & Isolir Masih Manual", "70%+ ISP lokal kehilangan 15–20% potensi pendapatan akibat keterlambatan isolir dan pencatatan tagihan di Excel / WhatsApp manual."),
        ("Fragmentasi Sistem & Perangkat", "Router MikroTik, OLT ZTE/Huawei, dan modem CPE diatur secara terpisah, membutuhkan waktu lama & biaya OPEX staf tinggi."),
        ("Buta Peta Jalur Lapangan", "Kabel fiber optic & ODP tidak terpetakan secara digital, memperlambat lokalisasi kabel putus dan penanganan gangguan."),
        ("Tingginya OPEX Operasional Staf", "Belum tersedianya portal mandiri untuk pelanggan, teknisi, & agen voucher membuat beban gaji karyawan membengkak.")
    ]

    for idx, (title, desc) in enumerate(problems):
        row = idx // 2
        col = idx % 2
        left = Inches(0.8 + col * 5.9)
        top = Inches(1.6 + row * 2.6)

        box = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(5.6), Inches(2.3))
        box.fill.solid()
        box.fill.fore_color.rgb = CARD_BG
        box.line.color.rgb = ACCENT_CYAN

        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.bold = True
        p.font.size = Pt(16)
        p.font.color.rgb = ACCENT_GOLD

        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(13)
        p2.font.color.rgb = TEXT_WHITE
        p2.space_before = Pt(8)

    # -------------------------------------------------------------
    # SLIDE 3: SOLUSI KITA (THE SOLUTION)
    # -------------------------------------------------------------
    slide3 = prs.slides.add_slide(blank_layout)
    add_background(slide3)
    add_header(slide3, "SLIDE 03", "Solusi Terpadu — Platform MyAdamedia")

    solutions = [
        ("Otomatisasi Billing & QRIS", "Integrasi QRIS Webhook instan. Tagihan dibayar, sistem otomatis mengubah status lunas & buka isolir MikroTik seketika."),
        ("Multi-Vendor Network Provisioning", "Mengontrol RouterOS MikroTik v6/v7 API, OLT PON via SNMP/Telnet, & TR-069 GenieACS CPE dalam satu dashboard."),
        ("Peta GIS Fiber & ODP Live", "Visualisasi Leaflet Satelit hybrid untuk lokasi pelanggan, titik ODP, & polyline rute kabel fiber optic."),
        ("Ekosistem 6 Multi-Portal", "Portal khusus Admin, Pelanggan Self-Service, Teknisi (GPS Ticket), Agen Voucher, Kolektor, & Investor.")
    ]

    for idx, (title, desc) in enumerate(solutions):
        row = idx // 2
        col = idx % 2
        left = Inches(0.8 + col * 5.9)
        top = Inches(1.6 + row * 2.6)

        box = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(5.6), Inches(2.3))
        box.fill.solid()
        box.fill.fore_color.rgb = CARD_BG
        box.line.color.rgb = ACCENT_CYAN

        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.bold = True
        p.font.size = Pt(16)
        p.font.color.rgb = ACCENT_CYAN

        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(13)
        p2.font.color.rgb = TEXT_WHITE
        p2.space_before = Pt(8)

    # -------------------------------------------------------------
    # SLIDE 4: TRAKTION & DATA OPERASIONAL AKTUAL
    # -------------------------------------------------------------
    slide4 = prs.slides.add_slide(blank_layout)
    add_background(slide4)
    add_header(slide4, "SLIDE 04", "Traksi & Kinerja Operasional Aktual")

    metrics = [
        ("80 Pelanggan", "Pelanggan Membayar Aktif\n(98,7% Active Rate)", ACCENT_CYAN),
        ("Rp 12.895.000", "Monthly Recurring Revenue\n(MRR Eksisting / Bulan)", ACCENT_GOLD),
        ("Rp 154.740.000", "Annualized Recurring Revenue\n(ARR Eksisting / Tahun)", ACCENT_CYAN),
        ("100% Fully Automation", "QRIS Auto-Release & Automatic\nIsolir MikroTik", TEXT_WHITE)
    ]

    for idx, (val, lbl, color) in enumerate(metrics):
        row = idx // 2
        col = idx % 2
        left = Inches(0.8 + col * 5.9)
        top = Inches(1.6 + row * 2.6)

        box = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(5.6), Inches(2.3))
        box.fill.solid()
        box.fill.fore_color.rgb = CARD_BG
        box.line.color.rgb = color

        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = val
        p.font.bold = True
        p.font.size = Pt(28)
        p.font.color.rgb = color

        p2 = tf.add_paragraph()
        p2.text = lbl
        p2.font.size = Pt(14)
        p2.font.color.rgb = TEXT_WHITE
        p2.space_before = Pt(10)

    # -------------------------------------------------------------
    # SLIDE 5: REVENUE BREAKDOWN PER PACKAGE
    # -------------------------------------------------------------
    slide5 = prs.slides.add_slide(blank_layout)
    add_background(slide5)
    add_header(slide5, "SLIDE 05", "Anatomi Revenue & Struktur Paket Pelanggan")

    pkgs = [
        ("Paket LITE", "Rp 150.000 / bln", "68 Pelanggan", "Rp 10.200.000 / bln (79,1%)"),
        ("Paket BASIC", "Rp 250.000 / bln", "7 Pelanggan", "Rp 1.750.000 / bln (13,6%)"),
        ("Paket BASIC A", "Rp 250.000 / bln", "3 Pelanggan", "Rp 750.000 / bln (5,8%)"),
        ("Paket STARTER A & B", "Rp 115.000 / bln", "3 Pelanggan", "Rp 345.000 / bln (1,5%)"),
    ]

    for idx, (pname, pprice, pcount, psum) in enumerate(pkgs):
        top = Inches(1.6 + idx * 1.3)
        box = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), top, Inches(11.733), Inches(1.1))
        box.fill.solid()
        box.fill.fore_color.rgb = CARD_BG
        box.line.color.rgb = ACCENT_CYAN

        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"{pname}  —  Harga: {pprice}  |  Jumlah: {pcount}  |  Pendapatan: {psum}"
        p.font.bold = True
        p.font.size = Pt(16)
        p.font.color.rgb = TEXT_WHITE

    # -------------------------------------------------------------
    # SLIDE 6: TECH STACK & QUALITY
    # -------------------------------------------------------------
    slide6 = prs.slides.add_slide(blank_layout)
    add_background(slide6)
    add_header(slide6, "SLIDE 06", "Arsitektur Kode & Kualitas Teknologi")

    tech_items = [
        ("Backend & Architecture", "Node.js v20+, Express.js, Clean Architecture, Repository Pattern, High Modular Design."),
        ("Database & Security", "SQLite (`better-sqlite3`), Encrypted Settings, Helmet Security Headers, Rate Limiting, Audit Logging."),
        ("Integration Engine", "MikroTik ROS v6/v7 Client, SNMP OLT ZTE/Huawei Engine, TR-069 ACS Server, Baileys WA Bot."),
        ("Production Reliability", "Jest Unit Test Suite (100% coverage pada modul kritis), Auto Daily Encrypted DB Backup.")
    ]

    for idx, (title, desc) in enumerate(tech_items):
        row = idx // 2
        col = idx % 2
        left = Inches(0.8 + col * 5.9)
        top = Inches(1.6 + row * 2.6)

        box = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(5.6), Inches(2.3))
        box.fill.solid()
        box.fill.fore_color.rgb = CARD_BG
        box.line.color.rgb = ACCENT_GOLD

        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.bold = True
        p.font.size = Pt(16)
        p.font.color.rgb = ACCENT_GOLD

        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(13)
        p2.font.color.rgb = TEXT_WHITE
        p2.space_before = Pt(8)

    # -------------------------------------------------------------
    # SLIDE 7: MARKET OPPORTUNITY
    # -------------------------------------------------------------
    slide7 = prs.slides.add_slide(blank_layout)
    add_background(slide7)
    add_header(slide7, "SLIDE 07", "Ukuran Pasar & Potensi Berskala")

    market_boxes = [
        ("TAM (Total Addressable Market)", "12.000+", "Pengelola RTRW Net & ISP Lokal di Indonesia"),
        ("SAM (Serviceable Addressable)", "3.500+", "ISP/RTRW Net di Pulau Jawa & Sumatera"),
        ("SOM Target (24 Bulan)", "250+", "Mitra ISP Terhubung Platform (SaaS MRR > Rp 100M)")
    ]

    for idx, (mtitle, mnum, mdesc) in enumerate(market_boxes):
        left = Inches(0.8 + idx * 3.9)
        box = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.8), Inches(3.7), Inches(4.5))
        box.fill.solid()
        box.fill.fore_color.rgb = CARD_BG
        box.line.color.rgb = ACCENT_CYAN

        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = mtitle
        p.font.bold = True
        p.font.size = Pt(14)
        p.font.color.rgb = ACCENT_CYAN

        p2 = tf.add_paragraph()
        p2.text = mnum
        p2.font.bold = True
        p2.font.size = Pt(36)
        p2.font.color.rgb = ACCENT_GOLD
        p2.space_before = Pt(20)

        p3 = tf.add_paragraph()
        p3.text = mdesc
        p3.font.size = Pt(13)
        p3.font.color.rgb = TEXT_WHITE
        p3.space_before = Pt(15)

    # -------------------------------------------------------------
    # SLIDE 8: BUSINESS MODEL
    # -------------------------------------------------------------
    slide8 = prs.slides.add_slide(blank_layout)
    add_background(slide8)
    add_header(slide8, "SLIDE 08", "Model Bisnis & Sumber Pendapatan")

    bmodels = [
        ("Direct ISP Subscription (B2C)", "Penjualan paket internet broadband bulanan langsung ke rumah tangga."),
        ("SaaS Licensing for ISPs (B2B)", "Model berlangganan software per ISP per bulan (Tiering Basic, Pro, Enterprise)."),
        ("Hotspot Voucher Revenue Share", "Komisi transaksi pembelian voucher hotspot publik via agen & QRIS mandiri.")
    ]

    for idx, (btitle, bdesc) in enumerate(bmodels):
        top = Inches(1.8 + idx * 1.7)
        box = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), top, Inches(11.733), Inches(1.4))
        box.fill.solid()
        box.fill.fore_color.rgb = CARD_BG
        box.line.color.rgb = ACCENT_CYAN

        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = btitle
        p.font.bold = True
        p.font.size = Pt(18)
        p.font.color.rgb = ACCENT_GOLD

        p2 = tf.add_paragraph()
        p2.text = bdesc
        p2.font.size = Pt(14)
        p2.font.color.rgb = TEXT_WHITE
        p2.space_before = Pt(8)

    # -------------------------------------------------------------
    # SLIDE 9: USE OF FUNDS (RP 20 JUTA)
    # -------------------------------------------------------------
    slide9 = prs.slides.add_slide(blank_layout)
    add_background(slide9)
    add_header(slide9, "SLIDE 09", "Kebutuhan Pendanaan & Alokasi Dana (Rp 20 Juta)")

    funds = [
        ("45%  |  Rp 9.000.000", "Hardware Expansion & Fiber Optic", "Pembelian OLT Splitter, Fiber Optic Core Cable, Enclosure Box, & Node ODP Baru."),
        ("30%  |  Rp 6.000.000", "Marketing & Penetrasi Pasar", "Digital Advertising, Spanduk, Banner Lokal, & Program Blast Pengenalan Layanan."),
        ("25%  |  Rp 5.000.000", "Server SSL, Cloud & OPEX Reserve", "Upgrade Server High Availability, Domain SSL, WhatsApp Official Gateway, & Cadangan OPEX.")
    ]

    for idx, (fval, ftitle, fdesc) in enumerate(funds):
        top = Inches(1.8 + idx * 1.7)
        box = slide9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), top, Inches(11.733), Inches(1.4))
        box.fill.solid()
        box.fill.fore_color.rgb = CARD_BG
        box.line.color.rgb = ACCENT_GOLD

        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"{fval}   —   {ftitle}"
        p.font.bold = True
        p.font.size = Pt(18)
        p.font.color.rgb = ACCENT_CYAN

        p2 = tf.add_paragraph()
        p2.text = fdesc
        p2.font.size = Pt(13)
        p2.font.color.rgb = TEXT_WHITE
        p2.space_before = Pt(6)

    # -------------------------------------------------------------
    # SLIDE 10: FINANCIAL PROJECTIONS (3 YEARS)
    # -------------------------------------------------------------
    slide10 = prs.slides.add_slide(blank_layout)
    add_background(slide10)
    add_header(slide10, "SLIDE 10", "Proyeksi Pertumbuhan Keuangan 3 Tahun")

    projs = [
        ("Tahun 1 (Target)", "180 Pelanggan", "MRR Rp 28.800.000", "Gross ARR Rp 345,6 Juta"),
        ("Tahun 2 (Target)", "380 Pelanggan", "MRR Rp 60.800.000", "Gross ARR Rp 729,6 Juta"),
        ("Tahun 3 (Target)", "750 Pelanggan", "MRR Rp 123.750.000", "Gross ARR Rp 1,48 Miliar")
    ]

    for idx, (pyear, pcust, pmrr, parr) in enumerate(projs):
        left = Inches(0.8 + idx * 3.9)
        box = slide10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.8), Inches(3.7), Inches(4.5))
        box.fill.solid()
        box.fill.fore_color.rgb = CARD_BG
        box.line.color.rgb = ACCENT_CYAN

        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = pyear
        p.font.bold = True
        p.font.size = Pt(16)
        p.font.color.rgb = ACCENT_GOLD

        p2 = tf.add_paragraph()
        p2.text = pcust
        p2.font.bold = True
        p2.font.size = Pt(22)
        p2.font.color.rgb = TEXT_WHITE
        p2.space_before = Pt(15)

        p3 = tf.add_paragraph()
        p3.text = pmrr
        p3.font.bold = True
        p3.font.size = Pt(18)
        p3.font.color.rgb = ACCENT_CYAN
        p3.space_before = Pt(10)

        p4 = tf.add_paragraph()
        p4.text = parr
        p4.font.size = Pt(14)
        p4.font.color.rgb = TEXT_MUTED
        p4.space_before = Pt(10)

    # -------------------------------------------------------------
    # SLIDE 11: INVESTOR RETURN & PAYBACK
    # -------------------------------------------------------------
    slide11 = prs.slides.add_slide(blank_layout)
    add_background(slide11)
    add_header(slide11, "SLIDE 11", "Estimasi ROI & Penawaran Investor")

    returns = [
        ("12% Profit Share", "Skema Bagi Hasil Bulanan dari Net Profit Bersih Bisnis."),
        ("5 – 7 Bulan", "Estimasi Masa Impas Modal (Payback Period) yang Sangat Cepat."),
        ("207% Net ROI", "Total Est. Pengembalian 24 Bulan: Rp 41.472.000 (Modal Rp 20 Juta).")
    ]

    for idx, (rtitle, rdesc) in enumerate(returns):
        left = Inches(0.8 + idx * 3.9)
        box = slide11.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.8), Inches(3.7), Inches(4.5))
        box.fill.solid()
        box.fill.fore_color.rgb = CARD_BG
        box.line.color.rgb = ACCENT_GOLD

        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = rtitle
        p.font.bold = True
        p.font.size = Pt(22)
        p.font.color.rgb = ACCENT_GOLD

        p2 = tf.add_paragraph()
        p2.text = rdesc
        p2.font.size = Pt(14)
        p2.font.color.rgb = TEXT_WHITE
        p2.space_before = Pt(20)

    # -------------------------------------------------------------
    # SLIDE 12: CALL TO ACTION & PENUTUP
    # -------------------------------------------------------------
    slide12 = prs.slides.add_slide(blank_layout)
    add_background(slide12)

    card12 = slide12.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(1.5), Inches(10.333), Inches(4.5))
    card12.fill.solid()
    card12.fill.fore_color.rgb = CARD_BG
    card12.line.color.rgb = ACCENT_CYAN
    card12.line.width = Pt(1.5)

    tx_cta = slide12.shapes.add_textbox(Inches(1.8), Inches(1.8), Inches(9.733), Inches(3.9))
    tf12 = tx_cta.text_frame
    tf12.word_wrap = True

    p = tf12.paragraphs[0]
    p.text = "Mari Bergabung Sebagai Mitra Strategis Kami!"
    p.font.bold = True
    p.font.size = Pt(32)
    p.font.color.rgb = ACCENT_CYAN

    p2 = tf12.add_paragraph()
    p2.text = "Bersama-sama Menguasai Pasar ISP Lokal & Digitalisasi Jaringan Komunitas"
    p2.font.size = Pt(18)
    p2.font.color.rgb = TEXT_WHITE
    p2.space_before = Pt(10)

    p3 = tf12.add_paragraph()
    p3.text = "Langkah Selanjutnya: Demo Platform Live  ➔  Penandatanganan Perjanjian  ➔  Eksekusi Ekspansi"
    p3.font.bold = True
    p3.font.size = Pt(14)
    p3.font.color.rgb = ACCENT_GOLD
    p3.space_before = Pt(20)

    p4 = tf12.add_paragraph()
    p4.text = "Management & Technical Team MyAdamedia\nEmail: admin@myadamedia.com  |  Portal: https://myadamedia.com"
    p4.font.size = Pt(13)
    p4.font.color.rgb = TEXT_MUTED
    p4.space_before = Pt(20)

    prs.save(pptx_path)
    print(f"PPTX Generated successfully at: {pptx_path}")

if __name__ == '__main__':
    pptx_path = r'd:\WEBAPP\myadamedia-billing\investor\PITCH_DECK.pptx'
    create_pitch_deck(pptx_path)
